import os
import fitz # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Convert SVG files to PNG using PyMuPDF
svg_dir = "data"
svg_files = [
    "spectrum_vs_angle.svg",
    "spectrum_orders.svg",
    "kbr_spectrum_vs_angle.svg",
    "kbr_spectrum_orders.svg",
    "peak_fit.svg",
    "collimator_comparison.svg",
    "collimator_comparison_raw.svg",
    "calibration_curves.svg"
]

png_map = {}
for svg in svg_files:
    svg_path = os.path.join(svg_dir, svg)
    png_path = os.path.join(svg_dir, svg.replace(".svg", ".png"))
    png_map[svg] = png_path
    if os.path.exists(svg_path):
        doc = fitz.open(svg_path)
        page = doc[0]
        pix = page.get_pixmap(dpi=300)
        pix.save(png_path)
        print(f"Converted {svg} -> {png_path}")

print("All SVGs converted successfully!")

# 2. Build PPTX presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
COLOR_BG = RGBColor(255, 255, 255)
COLOR_NAVY = RGBColor(15, 30, 60)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_MUTED = RGBColor(71, 85, 105)
COLOR_ACCENT = RGBColor(29, 78, 216)
COLOR_CARD_BG = RGBColor(248, 250, 252)
COLOR_CARD_BORDER = RGBColor(226, 232, 240)

TOTAL_MAIN_SLIDES = 14

def set_rtl(p, align=PP_ALIGN.RIGHT):
    """Sets Right-To-Left direction and alignment on PowerPoint paragraph."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    p.alignment = align

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="עקיפת בראג וספקטרוסקופיית קרני רנטגן"):
    top_bar = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
    tf = top_bar.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = category_text
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT
    p0.font.name = "Heebo"
    set_rtl(p0, PP_ALIGN.RIGHT)
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY
    p1.font.name = "Heebo"
    set_rtl(p1, PP_ALIGN.RIGHT)

def add_footer(slide, slide_num, total_slides=TOTAL_MAIN_SLIDES):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.4))
    tf = footer_box.text_frame
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if slide_num <= total_slides:
        p.text = f"מציג: יובל הירשמן | סמינר במעבדה 4מח | שקף {slide_num} / {total_slides}"
    else:
        p.text = f"מציג: יובל הירשמן | סמינר במעבדה 4מח | נספח {slide_num - total_slides}"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_MUTED
    p.font.name = "Heebo"
    set_rtl(p, PP_ALIGN.CENTER)

def add_card(slide, left, top, width, height, title="", bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    
    if title:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_right = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.font.name = "Heebo"
        set_rtl(p, PP_ALIGN.RIGHT)
    return shape

blank_layout = prs.slide_layouts[6]

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_ACCENT
line.line.fill.background()

tb = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "מעבדה לפיזיקה 4מח | הפקולטה לפיסיקה, הטכניון"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.font.name = "Heebo"
set_rtl(p, PP_ALIGN.RIGHT)

p2 = tf.add_paragraph()
p2.text = "מדידת ספקטרום הפליטה של שפופרת רנטגן, כיול זוויתי של גוניומטר וחקירת עקיפת בראג בגבישי LiF ו-KBr"
p2.font.size = Pt(26)
p2.font.bold = True
p2.font.color.rgb = COLOR_NAVY
p2.font.name = "Heebo"
p2.space_before = Pt(14)
set_rtl(p2, PP_ALIGN.RIGHT)

p3 = tf.add_paragraph()
p3.text = "עקיפת בראג וספקטרוסקופיית קרני רנטגן – מצגת סמינר מסכמת"
p3.font.size = Pt(16)
p3.font.color.rgb = COLOR_MUTED
p3.font.name = "Heebo"
p3.space_before = Pt(10)
set_rtl(p3, PP_ALIGN.RIGHT)

add_card(slide1, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.0))
tb_info = slide1.shapes.add_textbox(Inches(1.0), Inches(4.65), Inches(11.333), Inches(1.7))
tf_info = tb_info.text_frame
tf_info.word_wrap = True

items = [
    ("מציג:", " יובל הירשמן"),
    ("מנחה מדריכה:", " ד\"ר הדיל ח'מיס"),
    ("תאריך הגשה:", " 14 ביולי 2026"),
    ("מטרת הניסוי:", " זיהוי קווי פליטה אופייניים של Mo, כיול שגיאת נקודת האפס, וחקירת השפעת קבוע הסריג ורוחב הצמצם")
]
for label, val in items:
    p_inf = tf_info.add_paragraph() if tf_info.paragraphs[0].text else tf_info.paragraphs[0]
    run1 = p_inf.add_run()
    run1.text = label
    run1.font.bold = True
    run1.font.size = Pt(13)
    run1.font.color.rgb = COLOR_NAVY
    run1.font.name = "Heebo"
    
    run2 = p_inf.add_run()
    run2.text = val
    run2.font.size = Pt(13)
    run2.font.color.rgb = COLOR_TEXT
    run2.font.name = "Heebo"
    set_rtl(p_inf, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 2: Outline
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "מתווה ההצגה (Outline)")
add_footer(slide2, 2)

outline_items = [
    ("1. רקע תיאורטי", "קרינת רנטגן אופיינית ורציפה (בלימה), חוק בראג להתאבכות בגבישים, והמרת זוויות לאנרגיה"),
    ("2. מערכת המדידה", "מפרט מערכת PHYWE, אנודת מוליבדן (Mo), מונוכרומטור, גוניומטר, מונה גייגר וגבישי LiF / KBr"),
    ("3. תוצאות וניתוח LiF", "מיפוי קווי K_alpha ו-K_beta, כיול היסט האפס, ואימות חוק בראג בסדרי עקיפה n=1,2,3"),
    ("4. תוצאות וניתוח KBr", "תצפית בסדרי עקיפה מרובים (עד n=6) וזיהוי סף הבליעה העצמית של ברום (Br K-edge ב-13.47 keV)"),
    ("5. השפעת הצמצם וקריטריונים", "בחינת רזולוציה מול שטף (2mm מול 5mm), התאמה גאוסיאנית, והשוואת קריטריוני חילוץ אנרגיה")
]

for i, (head, desc) in enumerate(outline_items):
    top_pos = Inches(1.5 + i * 1.05)
    add_card(slide2, Inches(0.8), top_pos, Inches(11.733), Inches(0.9))
    tb_o = slide2.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.333), Inches(0.7))
    tf_o = tb_o.text_frame
    tf_o.word_wrap = True
    
    p_h = tf_o.paragraphs[0]
    p_h.text = head
    p_h.font.size = Pt(14)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_ACCENT
    p_h.font.name = "Heebo"
    set_rtl(p_h, PP_ALIGN.RIGHT)
    
    p_d = tf_o.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT
    p_d.font.name = "Heebo"
    set_rtl(p_d, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 3: Theoretical Background
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "עקרון עקיפת בראג וקרינת רנטגן אופיינית")
add_footer(slide3, 3)

add_card(slide3, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "עקרונות פיזיקליים מרכזיים")
tb_t3 = slide3.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(6.1), Inches(4.5))
tf_t3 = tb_t3.text_frame
tf_t3.word_wrap = True

bullets3 = [
    "• קרינת בלימה (Bremsstrahlung): ספקטרום רציף הנובע מהאטת אלקטרונים המואצים במתח 10-35 kV אל אנודת המוליבדן (Mo).",
    "• קרינה אופיינית: אלקטרונים מהירים מייננים אלקטרון מקליפה פנימית (K). מעבר אלקטרון מרמה גבוהה פולט פוטון באנרגיה בדידה: h ν = E_initial - E_final.",
    "• חוק בראג להתאבכות בונה: קרניים המוחזרות ממישורי גביש במרחק d יוצרות התאבכות בונה כאשר: 2d sin(θ_B) = n λ.",
    "• חילוץ אנרגיית הפוטון: קישור אורך הגל לאנרגיה מקבל את הצורה: E = (h c n) / (2 d sin θ_B).",
    "• קבועי הסריג שנבדקו: LiF (d = 201.4 pm) ו-KBr (d = 329.9 pm)."
]
for b in bullets3:
    p_b = tf_t3.add_paragraph() if tf_t3.paragraphs[0].text else tf_t3.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(12)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(8)
    set_rtl(p_b, PP_ALIGN.RIGHT)

add_card(slide3, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "איור: עקיפת בראג במישורי גביש")
if os.path.exists("images/bragg_angle.png"):
    slide3.shapes.add_picture("images/bragg_angle.png", Inches(7.7), Inches(2.2), width=Inches(4.633))

# ==========================================
# SLIDE 4: Experimental Setup
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "מערכת המדידה וסריקה זוויתית בגוניומטר")
add_footer(slide4, 4)

add_card(slide4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), "מערכת הניסוי (PHYWE X-Ray System)")
if os.path.exists("images/system1.png"):
    slide4.shapes.add_picture("images/system1.png", Inches(1.0), Inches(2.2), width=Inches(5.1))

add_card(slide4, Inches(6.5), Inches(1.5), Inches(6.033), Inches(5.3), "רכיבי המערכת ופרמטרי תפעול")
tb_t4 = slide4.shapes.add_textbox(Inches(6.7), Inches(2.1), Inches(5.633), Inches(4.5))
tf_t4 = tb_t4.text_frame
tf_t4.word_wrap = True

bullets4 = [
    "• שפופרת רנטגן: אנודת מוליבדן (Mo), מתח מאיץ עד 35 kV, זרם אלומה 1.0 mA.",
    "• צמצמים (Collimators): הלבשת דיאפרגרמה בקוטר 2mm (רזולוציה גבוהה) או 5mm (שטף גבוה) למיקוד אלומת הרנטגן.",
    "• מד זווית (גוניומטר): מאפשר סיבוב מסונכרן של הגביש בזווית θ וזרוע הגלאי בזווית 2θ לזיהוי זוויות בראג.",
    "• גלאי: מונה גייגר-מולר (Geiger-Müller counter) למדידת עוצמת הפיזור בקצב ספירה (counts per second - cps).",
    "• כיול זוויתי: ביצוע סריקת אפס מוקדמת לקביעת היסט האפס הכולל (Δθ_B) של הגוניומטר לפני כל סדרת מדידות."
]
for b in bullets4:
    p_b = tf_t4.add_paragraph() if tf_t4.paragraphs[0].text else tf_t4.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(12)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(10)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 5: Experimental Workflow
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "מהלך המדידה וסדרת הניסויים")
add_footer(slide5, 5)

cards5 = [
    ("מדידה 1: גביש LiF (2mm)", "סריקה זוויתית של גביש LiF (d = 201.4 pm) עם צמצם 2mm.\n\n• זיהוי קווי הפליטה K_alpha ו-K_beta בסדרים n=1, 2, 3.\n• חילוץ היסט אפס Δθ_LiF = +0.11° ± 0.06°.\n• בדיקת התאמה גאוסיאנית וכיול אנרגטי."),
    ("מדידה 2: גביש KBr (2mm)", "סריקה זוויתית של גביש KBr (d = 329.9 pm) עם צמצם 2mm.\n\n• קבוע סריג גדול יותר המניב זוויות קטנות.\n• תצפית בסדרי עקיפה מרובים (עד n=6).\n• גילוי סף הבליעה העצמית של ברום ב-13.47 keV."),
    ("מדידה 3: גביש LiF (5mm)", "סריקה זוויתית של גביש LiF עם צמצם 5mm.\n\n• השוואת השפעת רוחב הצמצם על הרזולוציה.\n• ניתוח ההתרחבות הגיאומטרית של השיאים (FWHM).\n• השוואת עוצמות השטף האבסולוטי (cps).")
]

for i, (title, content) in enumerate(cards5):
    left_pos = Inches(0.8 + i * 3.98)
    add_card(slide5, left_pos, Inches(1.5), Inches(3.78), Inches(5.3), title)
    tb_c5 = slide5.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.1), Inches(3.48), Inches(4.5))
    tf_c5 = tb_c5.text_frame
    tf_c5.word_wrap = True
    p_c5 = tf_c5.paragraphs[0]
    p_c5.text = content
    p_c5.font.size = Pt(12)
    p_c5.font.color.rgb = COLOR_TEXT
    p_c5.font.name = "Heebo"
    set_rtl(p_c5, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 6: LiF Spectrum
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "מיפוי קווי הפליטה K_alpha ו-K_beta בגביש LiF (צמצם 2mm)")
add_footer(slide6, 6)

add_card(slide6, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "ספקטרום עוצמה כפונקציה של זווית בראג (LiF)")
if os.path.exists(png_map["spectrum_vs_angle.svg"]):
    slide6.shapes.add_picture(png_map["spectrum_vs_angle.svg"], Inches(1.0), Inches(2.1), width=Inches(6.1))

add_card(slide6, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "ממצאים מרכזיים וזוויות שיא")
tb_t6 = slide6.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.633), Inches(4.5))
tf_t6 = tb_t6.text_frame
tf_t6.word_wrap = True

bullets6 = [
    "• היסט אפס כולל שנמצא: Δθ_B = +0.11° ± 0.06°.",
    "• קווי פליטה סדר 1: K_alpha ב-10.30° (SNR 55.1), K_beta ב-9.20° (SNR 23.5).",
    "• קווי פליטה סדר 2: K_alpha ב-20.70° (SNR 20.2), K_beta ב-18.30° (SNR 4.9).",
    "• קווי פליטה סדר 3: K_alpha ב-32.10° (SNR 6.7).",
    "• הפרדה ברורה: הצמצם הצר (2mm) מעניק הפרדה ספקטרלית מצוינת בין קו K_alpha לקו K_beta בכל סדרי הדיפרקציה."
]
for b in bullets6:
    p_b = tf_t6.add_paragraph() if tf_t6.paragraphs[0].text else tf_t6.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(12)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(8)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 7: Bragg's Law Verification across Orders (LiF)
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_header(slide7, "אימות חוק בראג עבור סדרי עקיפה שונים בגביש LiF")
add_footer(slide7, 7)

add_card(slide7, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "ספקטרום עוצמה כפונקציה של אנרגיה (n=1,2,3)")
if os.path.exists(png_map["spectrum_orders.svg"]):
    slide7.shapes.add_picture(png_map["spectrum_orders.svg"], Inches(1.0), Inches(2.1), width=Inches(6.1))

add_card(slide7, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "התכנסות לאנרגיה התיאורטית")
tb_t7 = slide7.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.633), Inches(4.5))
tf_t7 = tb_t7.text_frame
tf_t7.word_wrap = True

bullets7 = [
    "• המרת זוויות לאנרגיה: שימוש בחוק בראג עבור סדרים n=1, 2, 3.",
    "• הפחתת השפעת היסט האפס: ככל שסדר העקיפה n עולה, הזווית הנמדדת גדול יותר, ולכן השפעת היסט האפס הקבוע (Δθ) קטנה יחסית.",
    "• אנרגיות שיא מחושבות:",
    "  - סדר 1: K_a = 17.22 keV (3.2σ), K_b = 19.25 keV (3.4σ)",
    "  - סדר 2: K_a = 17.34 keV (3.6σ), K_b = 19.50 keV (2.1σ)",
    "  - סדר 3: K_a = 17.43 keV (2.2σ) – הקרוב ביותר לספרות (17.479 keV)!"
]
for b in bullets7:
    p_b = tf_t7.add_paragraph() if tf_t7.paragraphs[0].text else tf_t7.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(6)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 8: KBr Spectrum & Multiple Orders
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_header(slide8, "ספקטרום הפיזור בגביש KBr וסדרי עקיפה מרובים")
add_footer(slide8, 8)

add_card(slide8, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "ספקטרום עוצמה כפונקציה של זווית בראג (KBr)")
if os.path.exists(png_map["kbr_spectrum_vs_angle.svg"]):
    slide8.shapes.add_picture(png_map["kbr_spectrum_vs_angle.svg"], Inches(1.0), Inches(2.1), width=Inches(6.1))

add_card(slide8, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "תופעות סריג ייחודיות בגביש KBr")
tb_t8 = slide8.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.633), Inches(4.5))
tf_t8 = tb_t8.text_frame
tf_t8.word_wrap = True

bullets8 = [
    "• היסט אפס כולל: Δθ_KBr = +0.12° ± 0.25°.",
    "• זיהוי 8 שיאים בדידים ברורים עד לסדר עקיפה n=6.",
    "• השפעת קבוע הסריג: קבוע סריג גדול יותר (d_KBr = 329.9 pm לעומת d_LiF = 201.4 pm) יוצר זוויות פיזור קטנות יותר עבור אותה אנרגיה.",
    "• תוצאה: מאפשר ליותר סדרי פיזור להכנס לתחום הזוויות המוגבל של הגוניומטר (0°-40°).",
    "• רעש ועוצמה בסדרים גבוהים: בסדר 1 התקבלו התוצאות המדויקות ביותר (K_a = 17.40 keV, סטייה 0.6σ בלבד)."
]
for b in bullets8:
    p_b = tf_t8.add_paragraph() if tf_t8.paragraphs[0].text else tf_t8.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(7)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 9: Bromine K-edge Self-Absorption
# ==========================================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_header(slide9, "זיהוי סף הבליעה העצמית של ברום (Br K-edge) בגביש KBr")
add_footer(slide9, 9)

add_card(slide9, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "ספקטרום KBr לפי אנרגיה (הדגשת Br K-edge)")
if os.path.exists(png_map["kbr_spectrum_orders.svg"]):
    slide9.shapes.add_picture(png_map["kbr_spectrum_orders.svg"], Inches(1.0), Inches(2.1), width=Inches(6.1))

add_card(slide9, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "מנגנון פיזיקלי של בליעה עצמית")
tb_t9 = slide9.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.633), Inches(4.5))
tf_t9 = tb_t9.text_frame
tf_t9.word_wrap = True

bullets9 = [
    "• ירידה חדה בעוצמה ב-13.47 keV: בספקטרום המתקבל מגביש KBr נבחין בצניחה פתאומית וחדה של עוצמת קרינת הרקע.",
    "• התאמה לסף הבליעה K של ברום (Br): האנרגיה 13.47 keV מתאימה בדיוק לאנרגיית היינון של אלקטרון מקליפת K באטומי הברום.",
    "• מעבר לאל-אלסטיות: מעל אנרגיה זו, הפוטונים אינם עוברים רק פיזור אלסטי (בראג), אלא מייננים את אטומי הברום בגביש.",
    "• היעלמות פיקים: הבליעה העצמית החזקה מפחיתה את החזרתיות האלסטית ומסבירה את היעלמות פיק K_beta בסדר השני אל תוך הרעש."
]
for b in bullets9:
    p_b = tf_t9.add_paragraph() if tf_t9.paragraphs[0].text else tf_t9.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(7)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 10: Double Gaussian Fitting & Energy Calibration
# ==========================================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_header(slide10, "התאמה גאוסיאנית כפולה וכיול אנרגטי מדויק")
add_footer(slide10, 10)

add_card(slide10, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "התאמה גאוסיאנית לקווי K_alpha ו-K_beta (LiF 2mm)")
if os.path.exists(png_map["peak_fit.svg"]):
    slide10.shapes.add_picture(png_map["peak_fit.svg"], Inches(1.0), Inches(2.1), width=Inches(6.1))

add_card(slide10, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "פרמטרי ההתאמה והערכים שהתקבלו")
tb_t10 = slide10.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.633), Inches(4.5))
tf_t10 = tb_t10.text_frame
tf_t10.word_wrap = True

bullets10 = [
    "• מודל התאמה: גאוסיאן כפול עבור שני קווי הפליטה בתוספת רקע ליניארי למזעור רעשי הפיזור.",
    "• אנרגיית K_alpha שהתקבלה:",
    "  17.206 ± 0.083 keV (ספרות: 17.479 keV)",
    "• אנרגיית K_beta שהתקבלה:",
    "  19.338 ± 0.105 keV (ספרות: 19.608 keV)",
    "• הערכת אי-וודאות: השגיאות משקללות את שגיאת ההתאמה הסטטיסטית ואת רגישות הרזולוציה הזוויתית של הגוניומטר.",
    "• דיוק גבוה: התאמה רציפה מאפשרת חילוץ מרכזי שיא מעבר לצעד הדיסקרטי של המדידה (0.1°)."
]
for b in bullets10:
    p_b = tf_t10.add_paragraph() if tf_t10.paragraphs[0].text else tf_t10.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(6)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 11: Collimator Aperture Width Comparison
# ==========================================
slide11 = prs.slides.add_slide(blank_layout)
set_slide_background(slide11)
add_header(slide11, "השפעת רוחב הצמצם (2mm מול 5mm) על הרזולוציה והשטף")
add_footer(slide11, 11)

add_card(slide11, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.3), "ספקטרום מנורמל – הרחבת שיאים (FWHM)")
if os.path.exists(png_map["collimator_comparison.svg"]):
    slide11.shapes.add_picture(png_map["collimator_comparison.svg"], Inches(0.95), Inches(2.1), width=Inches(5.4))

add_card(slide11, Inches(6.833), Inches(1.5), Inches(5.7), Inches(5.3), "ספקטרום גולמי (cps) – גידול אקספוננציאלי בשטף")
if os.path.exists(png_map["collimator_comparison_raw.svg"]):
    slide11.shapes.add_picture(png_map["collimator_comparison_raw.svg"], Inches(6.983), Inches(2.1), width=Inches(5.4))

# ==========================================
# SLIDE 12: Energy Peak Extraction Criteria Comparison
# ==========================================
slide12 = prs.slides.add_slide(blank_layout)
set_slide_background(slide12)
add_header(slide12, "השוואת קריטריונים לחילוץ אנרגיות השיא")
add_footer(slide12, 12)

rows, cols = 7, 5
left_t, top_t, width_t, height_t = Inches(0.8), Inches(1.5), Inches(11.733), Inches(3.6)
table_shape = slide12.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
table = table_shape.table

headers = ["צמצם וקו קרינה", "עוצמה מרבית (Max Intensity)", "מרכז כובד (Centroid)", "התאמה גאוסיאנית (Gaussian)", "ערך ספרותי (keV)"]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Heebo"
    set_rtl(p, PP_ALIGN.CENTER)

table_data = [
    ["LiF 2mm - K_alpha", "17.215 ± 0.083 (0.3% | 3.2σ)", "17.204 ± 0.083 (0.3% | 3.3σ)", "17.206 ± 0.083 (0.3% | 3.3σ)", "17.479"],
    ["LiF 2mm - K_beta", "19.338 ± 0.105 (1.4% | 2.6σ)", "19.324 ± 0.105 (1.4% | 2.7σ)", "19.338 ± 0.105 (1.4% | 2.6σ)", "19.608"],
    ["KBr 2mm - K_alpha", "17.400 ± 0.134 (0.5% | 0.6σ)", "17.185 ± 0.134 (1.7% | 2.2σ)", "17.185 ± 0.134 (1.7% | 2.2σ)", "17.479"],
    ["KBr 2mm - K_beta", "19.610 ± 0.170 (0.0% | 0.0σ)", "19.290 ± 0.170 (1.6% | 1.9σ)", "19.290 ± 0.170 (1.6% | 1.9σ)", "19.608"],
    ["LiF 5mm - K_alpha", "17.215 ± 0.083 (0.3% | 3.2σ)", "17.430 ± 0.083 (0.3% | 0.6σ)", "17.464 ± 0.083 (0.1% | 0.2σ)", "17.479"],
    ["LiF 5mm - K_beta", "19.338 ± 0.105 (1.4% | 2.6σ)", "19.580 ± 0.105 (0.1% | 0.3σ)", "19.338 ± 0.105 (1.4% | 2.6σ)", "19.608"]
]

for i, row in enumerate(table_data):
    for j, val in enumerate(row):
        cell = table.cell(i+1, j)
        cell.text = val
        cell.fill.solid()
        if (i in [0, 2, 3] and j == 1) or (i == 4 and j == 3) or (i == 5 and j == 2):
            cell.fill.fore_color.rgb = RGBColor(220, 252, 231)
        else:
            cell.fill.fore_color.rgb = COLOR_CARD_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = "Heebo"
        set_rtl(p, PP_ALIGN.CENTER)

add_card(slide12, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.5), "מסקנת ניתוח הקריטריונים")
tb_b12 = slide12.shapes.add_textbox(Inches(1.0), Inches(5.7), Inches(11.333), Inches(1.0))
tf_b12 = tb_b12.text_frame
tf_b12.word_wrap = True
p_b12 = tf_b12.paragraphs[0]
p_b12.text = "• צמצם צר (2mm): קריטריון העוצמה המרבית (Max Intensity) מניב את הסטייה המינימלית מהערך הספרותי עקב פיקים חדים וברורים.\n• צמצם רחב (5mm): פיקים שטוחים ורחבים מביאים לכך שהתאמה גאוסיאנית ומרכז כובד מדויקים בהרבה, מכיוון שהם ממצעים את הרעש הסטטיסטי על פני כל הפרופיל."
p_b12.font.size = Pt(11.5)
p_b12.font.color.rgb = COLOR_TEXT
p_b12.font.name = "Heebo"
set_rtl(p_b12, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 13: Physical Discussion & Limitations
# ==========================================
slide13 = prs.slides.add_slide(blank_layout)
set_slide_background(slide13)
add_header(slide13, "דיון פיזיקלי במנגנוני הפיזור, הבליעה והמגבלות הניסיוניות")
add_footer(slide13, 13)

disc_items = [
    ("חשיבות כיול היסט האפס", "היסט זוויתי קטן (<0.5°) הנו קריטי לכיול האנרגטי, מכיוון ש-E ביחס הפוך ל-sin(θ). בסדרי עקיפה גבוהים השפעת ההיסט קטנה מורגשת פחות."),
    ("פשרת רוחב צמצם (Resolution vs Flux)", "רוחב הצמצם מכתיב פשרה הנדסית יסודית: צמצם רחב (5mm) משפר דרמטית את קצב הספירה (cps) אך מביא להרחבה גיאומטרית (FWHM) עקב התבדרות האלומה."),
    ("אינטראקציה לא-אלסטית (בליעה עצמית)", "תצפית ישירה בסף הבליעה K של ברום (13.47 keV) מראה את המעבר מפיזור בראג אלסטי ליינון קליפת K לא-אלסטי בגביש KBr."),
    ("השפעת קבוע הסריג (d)", "קבוע סריג גדול ב-KBr (329.9 pm לעומת 201.4 pm ב-LiF) דוחף את זוויות בראג לערכים קטנים יותר, ומאפשר צפייה בסדרי עקיפה מרובים (עד n=6).")
]

for i, (title, desc) in enumerate(disc_items):
    top_p = Inches(1.5 + i * 1.3)
    add_card(slide13, Inches(0.8), top_p, Inches(11.733), Inches(1.15), title)
    tb_d = slide13.shapes.add_textbox(Inches(1.0), top_p + Inches(0.4), Inches(11.333), Inches(0.7))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]
    p_d.text = desc
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT
    p_d.font.name = "Heebo"
    set_rtl(p_d, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 14: Conclusions
# ==========================================
slide14 = prs.slides.add_slide(blank_layout)
set_slide_background(slide14)
add_header(slide14, "מסקנות מרכזיות והישגי הניסוי")
add_footer(slide14, 14)

conc_cards = [
    ("אימות חוק בראג", "• שחזור מדויק של קווי הפליטה K_alpha ו-K_beta של מוליבדן בכל סדרי העקיפה.\n• התכנסות מיטבית לערך התיאורטי בסדרים גבוהים (סטייה של 0.05 keV בסדר n=3)."),
    ("כיול ורזולוציה", "• כיול היסט האפס של הגוניומטר חיוני להשגת דיוק אנרגטי תחרותי.\n• הבנה כמותית של הפשרה בין רזולוציה ספקטרלית (צמצם 2mm) לשטף קרינה (צמצם 5mm)."),
    ("תופעות סריג ובליעה", "• גילוי סף הבליעה העצמית K של ברום בגביש KBr ב-13.47 keV.\n• הדגמת השפעת קבוע הסריג d על טווח הזוויות הנמדד ומספר סדרי הדיפרקציה הנצפים."),
    ("אופטימיזציית ניתוח", "• במדידות חדות (צמצם 2mm) שיטת העוצמה המרבית עדיפה.\n• במדידות רחבות (צמצם 5mm) התאמה גאוסיאנית ומרכז כובד מספקות דיוק עדיף.")
]

for i, (title, content) in enumerate(conc_cards):
    row_i = i // 2
    col_i = i % 2
    left_p = Inches(0.8 + col_i * 5.98)
    top_p = Inches(1.5 + row_i * 2.7)
    add_card(slide14, left_p, top_p, Inches(5.75), Inches(2.5), title)
    tb_c = slide14.shapes.add_textbox(left_p + Inches(0.2), top_p + Inches(0.55), Inches(5.35), Inches(1.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = content
    p_c.font.size = Pt(12)
    p_c.font.color.rgb = COLOR_TEXT
    p_c.font.name = "Heebo"
    set_rtl(p_c, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 15: Appendix A - Calibration Curves
# ==========================================
slide15 = prs.slides.add_slide(blank_layout)
set_slide_background(slide15)
add_header(slide15, "נספח א': כיול היסט האפס של הגוניומטר (Δθ_B)", "נספחים להצגה")
add_footer(slide15, 15)

add_card(slide15, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3), "עקומות כיול זוויתי עבור LiF ו-KBr")
if os.path.exists(png_map["calibration_curves.svg"]):
    slide15.shapes.add_picture(png_map["calibration_curves.svg"], Inches(1.0), Inches(2.1), width=Inches(6.1))

add_card(slide15, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.3), "מתודולוגיית החישוב")
tb_a1 = slide15.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.633), Inches(4.5))
tf_a1 = tb_a1.text_frame
tf_a1.word_wrap = True

bullets_a1 = [
    "• שיטת התאמה: רגרסיה ליניארית בין הזוויות המדודות לתיאורטיות בכל השיאים.",
    "• היסט אפס לגביש LiF:",
    "  Δθ_LiF = +0.11° ± 0.06°",
    "• היסט אפס לגביש KBr:",
    "  Δθ_KBr = +0.12° ± 0.25°",
    "• משמעות פיזיקלית: שגיאה מכנית קבועה בזווית ההתחלתית של זרוע הגלאי בגוניומטר."
]
for b in bullets_a1:
    p_b = tf_a1.add_paragraph() if tf_a1.paragraphs[0].text else tf_a1.paragraphs[0]
    p_b.text = b
    p_b.font.size = Pt(12)
    p_b.font.color.rgb = COLOR_TEXT
    p_b.font.name = "Heebo"
    p_b.space_before = Pt(8)
    set_rtl(p_b, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 16: Appendix B - Full Data Table KBr
# ==========================================
slide16 = prs.slides.add_slide(blank_layout)
set_slide_background(slide16)
add_header(slide16, "נספח ב': טבלת נתוני שיאים מפורטת - גביש KBr", "נספחים להצגה")
add_footer(slide16, 16)

rows, cols = 9, 7
left_t, top_t, width_t, height_t = Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2)
table_shape16 = slide16.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
table16 = table_shape16.table

headers16 = ["קו ספקטרלי", "סדר (n)", "זווית תיאורטית (°)", "זווית מדודה (°)", "היסט (°)", "עוצמה (cps)", "SNR"]
for j, h in enumerate(headers16):
    cell = table16.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Heebo"
    set_rtl(p, PP_ALIGN.CENTER)

kbr_data = [
    ["K_alpha", "1", "6.17", "6.20", "+0.03", "773", "23.9"],
    ["K_alpha", "2", "12.42", "12.50", "+0.08", "427", "13.5"],
    ["K_alpha", "3", "18.82", "18.90", "+0.08", "164", "5.5"],
    ["K_alpha", "4", "25.47", "25.60", "+0.13", "119", "3.1"],
    ["K_alpha", "5", "32.52", "31.30", "-1.22", "113", "2.6"],
    ["K_beta", "1", "5.50", "5.50", "0.00", "268", "9.8"],
    ["K_beta", "3", "16.71", "17.40", "+0.69", "141", "3.5"],
    ["K_beta", "6", "35.10", "36.30", "+1.20", "125", "3.8"]
]

for i, row in enumerate(kbr_data):
    for j, val in enumerate(row):
        cell = table16.cell(i+1, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = "Heebo"
        set_rtl(p, PP_ALIGN.CENTER)

# Save presentation
output_pptx = "presentation.pptx"
prs.save(output_pptx)
print(f"RTL Presentation saved successfully to {output_pptx}!")
